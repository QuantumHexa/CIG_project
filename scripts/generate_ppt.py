"""Generate CIG Dev submission PowerPoint for EventLens."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent.parent / "docs" / "CIG_Project_Presentation.pptx"

# Brand colors
BG_DARK = RGBColor(15, 23, 42)       # slate-950
ACCENT = RGBColor(14, 165, 233)      # brand-500
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(148, 163, 184)
CARD = RGBColor(30, 41, 59)


def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CARD
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(9), Inches(0.4))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left=0.6, top=1.4, width=8.8, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)


def add_two_col(slide, left_items, right_items, left_title="", right_title=""):
    if left_title:
        lt = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.2), Inches(0.4))
        lp = lt.text_frame.paragraphs[0]
        lp.text = left_title
        lp.font.size = Pt(16)
        lp.font.bold = True
        lp.font.color.rgb = ACCENT
    if right_title:
        rt = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.2), Inches(0.4))
        rp = rt.text_frame.paragraphs[0]
        rp.text = right_title
        rp.font.size = Pt(16)
        rp.font.bold = True
        rp.font.color.rgb = ACCENT

    add_bullets(slide, left_items, left=0.6, top=1.75, width=4.2, height=4.8, size=16)
    add_bullets(slide, right_items, left=5.2, top=1.75, width=4.2, height=4.8, size=16)


def title_slide(prs, title, subtitle, extra=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(10), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.8))
    sp = s.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = ACCENT

    if extra:
        e = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(8.4), Inches(1.5))
        ep = e.text_frame.paragraphs[0]
        ep.text = extra
        ep.font.size = Pt(14)
        ep.font.color.rgb = MUTED


def content_slide(prs, title, bullets, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, title, subtitle)
    add_bullets(slide, bullets)


def table_slide(prs, title, headers, rows, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, title, subtitle)

    cols, rcount = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(rcount, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.45 * rcount)).table

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = ACCENT

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(
        prs,
        "EventLens",
        "Event & Media Management Platform",
        "CIG Dev Problem Statement\nTeam: QuantumHexa\nGitHub: github.com/QuantumHexa/CIG_project",
    )

    content_slide(
        prs,
        "The Problem",
        [
            "Club events generate hundreds of photos & videos scattered across:",
            "  • Google Drives, personal folders, WhatsApp groups",
            "  • No central place to organize, search, or share media",
            "",
            "Pain points:",
            "  • Hard to find your own photos in large albums",
            "  • No access control for private club content",
            "  • No smart search by tags, events, or people",
            "  • Manual watermarking and sharing is tedious",
        ],
    )

    content_slide(
        prs,
        "Our Solution — EventLens",
        [
            "A centralized, scalable web platform for clubs & photographers",
            "",
            "One place to:",
            "  • Create events & albums with metadata",
            "  • Bulk upload photos/videos with drag-and-drop",
            "  • Control public vs private access by role",
            "  • Interact socially — like, comment, tag, favorite",
            "  • AI auto-tagging + advanced search",
            "  • Find your photos via face-matching selfie",
            "  • Cloud storage (AWS S3) + watermarked downloads",
        ],
        "Production-ready full-stack application",
    )

    content_slide(
        prs,
        "Core Features — Event Management",
        [
            "Create & manage events with descriptions & metadata",
            "Event-wise albums (Main Stage, Backstage, etc.)",
            "Sort & filter by: Event name | Date | Category",
            "Club name & category tagging for organization",
            "Public / private event visibility",
        ],
    )

    content_slide(
        prs,
        "Core Features — Media Upload",
        [
            "Upload photos & videos — single or bulk (up to 20 files)",
            "Drag-and-drop upload zone with live preview",
            "Sharp compression & resize (max 2400px, JPEG 85%)",
            "Automatic thumbnail generation (400×400)",
            "Per-upload public / private toggle",
            "Windows-friendly: accepts JPG, PNG, WebP, HEIC, MP4",
        ],
    )

    content_slide(
        prs,
        "Authentication & Access Control",
        [
            "JWT-based secure authentication",
            "",
            "Four roles:",
            "  • Admin — full platform control",
            "  • Photographer — upload & manage events",
            "  • Club Member — access private club media",
            "  • Viewer — browse public content",
            "",
            "Public media: visible to everyone",
            "Private media: restricted to authorized members",
        ],
    )

    content_slide(
        prs,
        "Social Features & Notifications",
        [
            "Like, comment, favorite on any media item",
            "Tag friends/users in photos",
            "Share event albums via URL / QR endpoint",
            "Download with one click (authenticated)",
            "",
            "Real-time notifications (Socket.io):",
            "  • Someone liked your photo",
            "  • Someone commented on your upload",
            "  • Someone tagged you in a photo",
        ],
    )

    content_slide(
        prs,
        "AI / ML Features",
        [
            "Smart Image Tagging (on upload):",
            "  • Category-based tags (cultural, sports, crowd…)",
            "  • Color analysis (outdoor, nature, night…)",
            "  • Optional Hugging Face BLIP captions (HF_API_TOKEN)",
            "",
            "Advanced Search:",
            "  • Event name, tags, upload date, uploader name",
            "",
            "Facial Recognition / My Photos:",
            "  1. Upload reference selfie",
            "  2. Match against all event uploads",
            "  3. Personalized gallery of your photos",
        ],
    )

    content_slide(
        prs,
        "Cloud Integration & Watermarking",
        [
            "AWS S3 integration (production):",
            "  • PutObject upload, signed URLs",
            "  • Automatic fallback to local storage in dev",
            "",
            "Dynamic watermark on download:",
            "  • Club name + Event name + User role",
            "  • Applied via Sharp SVG composite",
            "  • Protects photographer & club branding",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Innovation & Bonus Features", "Beyond core requirements")
    add_two_col(
        slide,
        [
            "Infinite scroll gallery",
            "PWA manifest (installable app)",
            "QR-based album share API",
            "Docker Compose for PostgreSQL",
            "Full API documentation",
        ],
        [
            "Modular TypeScript backend",
            "TanStack Query on frontend",
            "Role-based route guards",
            "Seed data for instant demo",
            "Cross-platform face matching",
        ],
        "Bonus Features",
        "Engineering Quality",
    )

    content_slide(
        prs,
        "Technology Stack",
        [
            "Frontend:  React 19 + Vite + Tailwind CSS + TanStack Query",
            "Backend:   Node.js + Express + TypeScript",
            "Database:  Prisma ORM + SQLite (dev) / PostgreSQL (prod)",
            "Storage:   AWS S3 SDK v3 + Sharp image processing",
            "Realtime:  Socket.io (notifications)",
            "Auth:      JWT + bcrypt + role middleware",
            "Deploy:    Render/Railway (API) + Vercel (frontend)",
        ],
    )

    content_slide(
        prs,
        "System Architecture",
        [
            "Client (React SPA / PWA)",
            "       ↓ REST API + WebSocket",
            "Express Server",
            "  ├── Auth Middleware (JWT + Roles)",
            "  ├── REST Routes (events, media, users, notifications)",
            "  ├── AI Services (tagging, face embeddings)",
            "  ├── Watermark Service (Sharp)",
            "  └── Storage Service → AWS S3 / local uploads",
            "       ↓",
            "PostgreSQL / SQLite via Prisma ORM",
            "",
            "See docs/ARCHITECTURE.md for full Mermaid diagram",
        ],
    )

    content_slide(
        prs,
        "Database Schema (High Level)",
        [
            "User — roles, face descriptor for matching",
            "Event — name, date, category, club, public flag",
            "Album — event-wise album grouping",
            "Media — url, thumbnail, tags, face data, public flag",
            "MediaTag — AI / user labels for search",
            "Like, Comment, Favorite, MediaUserTag",
            "Notification — real-time alert records",
            "",
            "Full ERD: docs/DATABASE_SCHEMA.md + prisma/schema.prisma",
        ],
    )

    content_slide(
        prs,
        "Live Demo Flow",
        [
            "1. Sign in → photo@cig.dev / password123",
            "2. Events → Annual Cultural Fest 2026",
            "3. Upload → drag-drop 2–3 photos",
            "4. Gallery → view AI-generated tags",
            "5. Search → filter by tag (e.g. crowd)",
            "6. Like + comment → real-time notification",
            "7. My Photos → upload selfie → find matches",
            "8. Download → watermarked image (club + event + role)",
            "",
            "Local: localhost:5173  |  API: localhost:4000",
        ],
        "5–10 minute demo script for judges",
    )

    table_slide(
        prs,
        "Evaluation Criteria Alignment",
        ["Criteria", "Weight", "Our Implementation"],
        [
            ("UI/UX & Design", "15%", "Modern dark UI, drag-drop, responsive gallery"),
            ("Backend & APIs", "15%", "REST modules, Prisma, typed services"),
            ("Auth & Access", "10%", "JWT + 4 roles, public/private media"),
            ("Cloud Integration", "15%", "AWS S3 SDK, signed URLs, prod-ready"),
            ("Media Management", "15%", "Events, albums, bulk upload, compression"),
            ("AI/ML Features", "15%", "Auto tags, search, face match pipeline"),
            ("Notifications", "5%", "Socket.io real-time alerts"),
            ("Code Quality", "5%", "TypeScript, modular architecture"),
            ("Innovation", "5%", "PWA, infinite scroll, QR share"),
        ],
        "Mapped to CIG Dev marking rubric",
    )

    content_slide(
        prs,
        "Deliverables & Submission",
        [
            "✓ GitHub: github.com/QuantumHexa/CIG_project",
            "✓ README + setup instructions",
            "✓ Database schema documentation",
            "✓ Architecture diagram (Mermaid)",
            "✓ API documentation (docs/API.md)",
            "✓ This presentation",
            "○ Live deployed demo (Render + Vercel) — in progress",
            "○ Demo video (5–10 min walkthrough) — record next",
        ],
    )

    title_slide(
        prs,
        "Thank You",
        "Questions?",
        "EventLens — Centralize. Organize. Discover.\n\nGitHub: github.com/QuantumHexa/CIG_project",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Created: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
